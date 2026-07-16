"""
Generate executive PowerPoint for AI Analyst Bot evaluation (production run).
Output: AI_Analyst_Bot_Evaluation_Report.pptx in repo root.
"""
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT_PATH = ROOT / "AI_Analyst_Bot_Evaluation_Report.pptx"

# Theme
NAVY = RGBColor(0x1A, 0x36, 0x5D)
ACCENT = RGBColor(0x25, 0x63, 0xEB)
DARK = RGBColor(0x1E, 0x29, 0x3B)
MUTED = RGBColor(0x64, 0x74, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x05, 0x96, 0x69)
AMBER = RGBColor(0xD9, 0x77, 0x06)
RED = RGBColor(0xDC, 0x26, 0x26)


def set_slide_bg(slide, rgb: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def add_title_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(8.8), Inches(2.2))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AI Analyst Bot"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = "Full Evaluation Analytics — DTC Fleet Health Assistant"
    p2.font.size = Pt(22)
    p2.font.color.rgb = RGBColor(0xBF, 0xDB, 0xFE)
    p2.space_before = Pt(12)
    meta = slide.shapes.add_textbox(Inches(0.6), Inches(5.2), Inches(8.8), Inches(1.2))
    mtf = meta.text_frame
    mtf.paragraphs[0].text = (
        "Production benchmark · 1,000 sessions · 895 turns · 11 customers\n"
        "Run: eval_20260520_085108 · May 2026 · gpt-3.5-turbo · LLM judge"
    )
    for p in mtf.paragraphs:
        p.font.size = Pt(14)
        p.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)


def add_section_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, ACCENT)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.8), Inches(8.8), Inches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = RGBColor(0xDB, 0xEA, 0xFE)
        p2.space_before = Pt(10)


def add_bullet_slide(
    prs: Presentation,
    title: str,
    bullets: list[str],
    footer: str = "",
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    hdr = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9.0), Inches(0.7))
    hp = hdr.text_frame.paragraphs[0]
    hp.text = title
    hp.font.size = Pt(28)
    hp.font.bold = True
    hp.font.color.rgb = NAVY
    body = slide.shapes.add_textbox(Inches(0.55), Inches(1.15), Inches(9.0), Inches(5.5))
    tf = body.text_frame
    tf.word_wrap = True
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(16)
        p.font.color.rgb = DARK
        p.level = 0
        p.space_after = Pt(8)
    if footer:
        ft = slide.shapes.add_textbox(Inches(0.5), Inches(6.85), Inches(9.0), Inches(0.4))
        fp = ft.text_frame.paragraphs[0]
        fp.text = footer
        fp.font.size = Pt(10)
        fp.font.color.rgb = MUTED


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    col_widths: list[float] | None = None,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    hdr = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9.0), Inches(0.6))
    hdr.text_frame.paragraphs[0].text = title
    hdr.text_frame.paragraphs[0].font.size = Pt(24)
    hdr.text_frame.paragraphs[0].font.bold = True
    hdr.text_frame.paragraphs[0].font.color.rgb = NAVY

    nrows = len(rows) + 1
    ncols = len(headers)
    left, top, width, height = Inches(0.4), Inches(1.05), Inches(9.2), Inches(0.35 * nrows)
    table = slide.shapes.add_table(nrows, ncols, left, top, width, height).table

    if col_widths:
        for ci, w in enumerate(col_widths):
            table.columns[ci].width = Inches(w)

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF1, 0xF5, 0xF9)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = DARK


def add_kpi_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    hdr = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9.0), Inches(0.6))
    hdr.text_frame.paragraphs[0].text = "Executive scorecard"
    hdr.text_frame.paragraphs[0].font.size = Pt(28)
    hdr.text_frame.paragraphs[0].font.bold = True
    hdr.text_frame.paragraphs[0].font.color.rgb = NAVY

    kpis = [
        ("83.5", "AI health score", AMBER),
        ("91.4%", "Answer quality", GREEN),
        ("52.9%", "Session pass rate", AMBER),
        ("895", "Completed turns", DARK),
        ("35.3 s", "Avg latency", AMBER),
        ("26.4 s", "Median latency", DARK),
        ("$6.99", "Total API cost", DARK),
        ("0", "Hallucinations", GREEN),
    ]
    cols, rows = 4, 2
    for i, (val, label, color) in enumerate(kpis):
        r, c = divmod(i, cols)
        x = Inches(0.45 + c * 2.35)
        y = Inches(1.2 + r * 2.6)
        box = slide.shapes.add_shape(1, x, y, Inches(2.15), Inches(2.2))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
        box.line.color.rgb = RGBColor(0xE2, 0xE8, 0xF0)
        tf = box.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p0 = tf.paragraphs[0]
        p0.text = val
        p0.font.size = Pt(32)
        p0.font.bold = True
        p0.font.color.rgb = color
        p0.alignment = PP_ALIGN.CENTER
        p1 = tf.add_paragraph()
        p1.text = label
        p1.font.size = Pt(12)
        p1.font.color.rgb = MUTED
        p1.alignment = PP_ALIGN.CENTER

    callout = slide.shapes.add_textbox(Inches(0.5), Inches(6.0), Inches(9.0), Inches(1.0))
    ctf = callout.text_frame
    ctf.word_wrap = True
    ctf.paragraphs[0].text = (
        "Verdict: Strong when data is available; ~half of sessions fail strict gates due to "
        "SQL gaps and multi-turn timeouts. Zero safety or hallucination issues detected."
    )
    ctf.paragraphs[0].font.size = Pt(14)
    ctf.paragraphs[0].font.color.rgb = DARK


def add_chart_bar_slide(
    prs: Presentation,
    title: str,
    categories: list[str],
    values: list[float],
    x_label: str,
    note: str,
) -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    hdr = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9.0), Inches(0.6))
    hdr.text_frame.paragraphs[0].text = title
    hdr.text_frame.paragraphs[0].font.size = Pt(24)
    hdr.text_frame.paragraphs[0].font.bold = True
    hdr.text_frame.paragraphs[0].font.color.rgb = NAVY

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series("Score (%)", values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.5),
        Inches(1.0),
        Inches(9.0),
        Inches(5.0),
        chart_data,
    ).chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.gap_width = 80
    if note:
        nb = slide.shapes.add_textbox(Inches(0.5), Inches(6.2), Inches(9.0), Inches(0.5))
        nb.text_frame.paragraphs[0].text = note
        nb.text_frame.paragraphs[0].font.size = Pt(10)
        nb.text_frame.paragraphs[0].font.color.rgb = MUTED


def add_closing_slide(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, NAVY)
    box = slide.shapes.add_textbox(Inches(0.6), Inches(2.5), Inches(8.8), Inches(2.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Thank you"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p2 = tf.add_paragraph()
    p2.text = "Source: eval_20260520_085108_190cf6 · combined_report.json"
    p2.font.size = Pt(14)
    p2.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    p2.space_before = Pt(16)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)

    add_bullet_slide(
        prs,
        "Executive summary (plain English)",
        [
            "We tested the fleet AI assistant with 1,000 real business questions on production.",
            "When it gets data, answers score ~91% on quality — maintenance prioritization is near perfect (99%).",
            "About 53% of sessions pass all strict gates; failures are usually missing SQL data, not made-up facts.",
            "Average wait time is 35 seconds per question; follow-up (multi-turn) sessions average 11+ minutes.",
            "Total test cost was about $7 in API fees — roughly $0.007 per question.",
            "No hallucinations, no cross-customer data leaks, no unsafe SQL in 1,000 tests.",
        ],
    )

    add_kpi_slide(prs)

    add_section_slide(prs, "Quality", "By question type")

    quality_rows = [
        ["Maintenance Prioritization", "99.3%", "Which vehicles to service first — excellent"],
        ["Fleet DTC Distribution", "96.9%", "Fault spread across fleet — excellent"],
        ["Fleet Health", "96.5%", "Overall fleet snapshot — excellent"],
        ["Fleet Trends", "89.8%", "Fault trends over time — good"],
        ["Vehicle Investigation", "88.1%", "One vehicle deep-dive — good but flaky"],
        ["DTC Investigation", "85.8%", "One fault code deep-dive — good but flaky"],
        ["Adversarial / Edge Cases", "83.5%", "Trick questions — acceptable"],
        ["Co-occurrence Analysis", "82.8%", "Faults appearing together — weakest core feature"],
        ["Fleet Investigation", "N/A", "Multi-step queries often timed out"],
    ]
    add_table_slide(
        prs,
        "Quality by category (trace judge %)",
        ["Category", "Score", "Meaning"],
        quality_rows,
        col_widths=[3.2, 0.9, 5.1],
    )

    add_chart_bar_slide(
        prs,
        "Quality scores (higher is better)",
        [
            "Maint.",
            "DTC Dist.",
            "Fleet Hlth",
            "Trends",
            "Vehicle",
            "DTC Inv.",
            "Adversarial",
            "Co-occur.",
        ],
        [99.3, 96.9, 96.5, 89.8, 88.1, 85.8, 83.5, 82.8],
        "Quality (%)",
        "Production run · Maintenance and fleet summaries strongest · Co-occurrence weakest",
    )

    add_section_slide(prs, "Speed", "Latency and responsiveness")

    add_bullet_slide(
        prs,
        "Speed analysis",
        [
            "Single-turn (one question): mean 32.4 s · median 26.4 s · p95 72.7 s",
            "Multi-turn (follow-ups): mean 685 s (~11 min) · p95 1,918 s (~32 min)",
            "883 of 895 sessions finished under 109 seconds",
            "One session ran 2,156 seconds (36 minutes) — production risk",
            "SQL execution itself is fast (~0.37 s avg in dev traces); slowness is LLM + multiple tools",
            "Recommendation: 60–90 s hard timeout with partial results for chat UX",
        ],
        footer="Latency by session type · production benchmark",
    )

    add_table_slide(
        prs,
        "Latency distribution (session count by bucket)",
        ["Latency bucket", "Sessions"],
        [
            ["1–109 s", "883"],
            ["109–217 s", "6"],
            ["217–325 s", "3"],
            ["432–540 s", "1"],
            ["540–648 s", "1"],
            ["2,048–2,156 s", "1"],
        ],
        col_widths=[4.5, 4.7],
    )

    add_section_slide(prs, "Cost", "API spend and efficiency")

    add_table_slide(
        prs,
        "Cost by category (USD)",
        ["Category", "Cost (USD)"],
        [
            ["DTC Investigation", "$1.65"],
            ["Vehicle Investigation", "$1.38"],
            ["Fleet DTC Distribution", "$1.19"],
            ["Co-occurrence Analysis", "$1.11"],
            ["Maintenance Prioritization", "$0.84"],
            ["Fleet Trends", "$0.43"],
            ["Fleet Health", "$0.39"],
        ],
        col_widths=[5.5, 3.7],
    )

    add_bullet_slide(
        prs,
        "Cost insights",
        [
            "Total run: $6.99 for 1,000 sessions (~$0.007 per session)",
            "VRL Logistics: $6.21 (89% of spend) — largest fleet in test set",
            "~12,000 tokens per turn on average",
            "Co-occurrence: high cost + lowest quality = worst value",
            "Scale estimate: ~10,000 questions/month ≈ $70 API cost (model only)",
        ],
    )

    add_section_slide(prs, "Reliability", "SQL, tools, and failures")

    add_table_slide(
        prs,
        "SQL outcome distribution (895 sessions)",
        ["Outcome", "Sessions", "Share"],
        [
            ["High quality (score > 80%)", "820", "91.6%"],
            ["Very poor (score < 20%)", "29", "3.2%"],
            ["Mediocre (40–60%)", "29", "3.2%"],
            ["Fair (60–80%)", "9", "1.0%"],
            ["Low (20–40%)", "8", "0.9%"],
        ],
        col_widths=[4.0, 2.5, 2.7],
    )

    add_table_slide(
        prs,
        "Top tool usage",
        ["Tool", "Calls"],
        [
            ["get_maintenance_priority", "267"],
            ["get_dtc_cooccurrence", "251"],
            ["get_fleet_dtc_distribution", "207"],
            ["run_sql", "204"],
            ["get_fleet_trends", "123"],
            ["get_dtc_details", "109"],
        ],
        col_widths=[6.0, 3.2],
    )

    add_bullet_slide(
        prs,
        "Common failure patterns",
        [
            "insufficient_evidence — SQL returned 0 rows; bot says “could not fetch enough data”",
            "sql_error — malformed or failed query (common on vehicle/DTC lookups)",
            "exception:timed_out — rare but critical (up to 36 minutes)",
            "heuristic_fallback_used — safe fallback when data missing",
        ],
    )

    add_section_slide(prs, "Test evolution", "Three evaluation runs")

    add_table_slide(
        prs,
        "Run progression (pass rate context)",
        ["Run", "Scope", "Pass rate", "Health", "Notes"],
        [
            ["Run 1", "20 conversational", "5%", "68.4", "Dev laptop debugging"],
            ["Run 2", "50 single-turn", "80%", "90.7", "Promising smaller slice"],
            ["Run 3", "1,000 production", "52.9%", "83.5", "Full benchmark; harder scenarios"],
        ],
        col_widths=[1.2, 2.2, 1.2, 1.0, 4.6],
    )

    add_bullet_slide(
        prs,
        "Why pass rate dropped 80% → 53%",
        [
            "Run 3 is not “worse code” — it is a harder, larger test.",
            "Adds multi-turn sessions (mean 685 s latency, many failures).",
            "Adds co-occurrence, fleet investigation, 11 customers.",
            "Stricter session-level pass gates vs. smaller dev samples.",
        ],
    )

    add_section_slide(prs, "Failures", "Sample worst sessions")

    add_table_slide(
        prs,
        "Worst sessions (sample)",
        ["Session", "Category", "Trace", "Batch", "Latency"],
        [
            ["sess_single_veh_16", "Vehicle Inv.", "26%", "0%", "33 s"],
            ["sess_single_veh_13", "Vehicle Inv.", "47%", "0%", "34 s"],
            ["sess_single_gen_0203", "Vehicle Inv.", "15%", "20%", "242 s"],
            ["sess_single_gen_0124", "DTC Inv.", "35%", "0%", "131 s"],
            ["sess_single_gen_0129", "Co-occurrence", "70%", "0%", "164 s"],
            ["sess_single_dtc_03", "DTC Inv.", "65%", "0%", "30 s"],
        ],
        col_widths=[2.8, 2.0, 1.2, 1.2, 1.0],
    )

    add_section_slide(prs, "Recommendations", "Prioritized actions")

    add_bullet_slide(
        prs,
        "Top 5 fixes (by impact)",
        [
            "1. Multi-turn timeouts — cap at 60–90 s; return partial results (mean 685 s today).",
            "2. Co-occurrence SQL — fix query templates / date ranges (lowest quality category).",
            "3. Vehicle/DTC lookup SQL — many worst sessions are narrow vehicle filters.",
            "4. Speed — cache fleet health and maintenance priority for common questions.",
            "5. Keep safety guardrails — zero hallucinations; do not trade safety for speed.",
        ],
    )

    add_table_slide(
        prs,
        "Report card — all dimensions",
        ["Dimension", "Grade", "Summary"],
        [
            ["Answer quality (when data exists)", "A", "~91% judge score"],
            ["Reliability (complete success)", "C+", "~53% pass rate"],
            ["Speed (single-turn)", "C+", "~35 s average"],
            ["Multi-turn support", "F", "Not production-ready"],
            ["Safety / hallucination", "A+", "Zero issues detected"],
            ["Cost efficiency", "B", "~$7 / 1k questions"],
            ["Scalability", "B−", "OK single-turn; stress on multi-turn"],
            ["Maintenance prioritization", "A+", "99.3% — ship with confidence"],
        ],
        col_widths=[3.5, 0.8, 5.0],
    )

    add_closing_slide(prs)

    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
